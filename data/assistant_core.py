# assistant_core.py
# Команды: только commands.parse_user / handle_user.
# Чат: LLM + персона. Эмоции: AnimationSelector + emotion_model.
# IntentRouter отключён — не исполняет и не подсказывает LLM.
import asyncio
import aiohttp
import json
import re
import os
from typing import Optional, Callable, AsyncGenerator, Dict, List, Any

from rag_engine import RAGEngine
from system_controller import SystemController
from voice_controller import VoiceController
from utils import task_pool, logger, api_semaphore

import config
from memory_clean import clean_reply, is_system_or_ocr, log_system
try:
    from prompt_builder import build_system_prompt
    HAS_PROMPT_BUILDER = True
except ImportError:
    HAS_PROMPT_BUILDER = False
    build_system_prompt = None
try:
    from animation_selector import AnimationSelector
    HAS_ANIM_SELECTOR = True
except ImportError:
    HAS_ANIM_SELECTOR = False
    AnimationSelector = None
from memory_manager import MemoryManager
from persistent_memory import PersistentMemory
from command_parser import CommandParser
from command_executor import CommandExecutor
from notes_manager import NotesManager
from context_manager import ContextManager
try:
    from memory_scope import (
        active_character,
        character_scope,
        prompt_scopes,
        scope_for_category,
    )
except ImportError:
    def active_character():
        return getattr(config, "ACTIVE_CHARACTER", None) or "лисичка"

    def character_scope(name=None):
        return f"character:{(name or active_character())}"

    def prompt_scopes(name=None):
        return ["user", "pc", "global", character_scope(name)]

    def scope_for_category(category):
        return "global"

# intent_router.py не подключаем: один вход команд — commands.py

# Микро-модели — только эмоции / морда, не маршрутизация команд
try:
    from micro_models import HybridAnalyzer
    MICRO_MODELS_AVAILABLE = True
except ImportError:
    MICRO_MODELS_AVAILABLE = False
    print("⚠️ micro_models.py не найден, используем старый анализатор")


class LMAssistant:
    """
    Асинхронное ядро ИИ-ассистента Лисичка.

    Команды — только commands.handle_user.
    Нет команды — чат (персона + RAG + память).
    Эмоции — AnimationSelector / emotion_model.
    Настроение — ContextManager.
    """

    def __init__(self):
        self.api_url = config.API_URL
        self.api_key = config.API_KEY
        self.model_name = config.MODEL_NAME
        self.temperature = config.TEMPERATURE
        self.max_tokens = config.MAX_TOKENS
        self.system_prompt = config.SYSTEM_PROMPT

        # Компоненты
        self.rag = RAGEngine(
            db_path=getattr(config, "ADVANCED_RAG_DB", "advanced_rag.db"),
            api_url=self.api_url,
            # embeddings-модель, НЕ chat (иначе FAISS мусор)
            model_name=getattr(config, "EMBEDDING_MODEL", None),
            dimension=int(getattr(config, "ADVANCED_RAG_DIMENSION", 768) or 768),
        )
        self.system = SystemController()
        self.voice = VoiceController()
        self.executor = CommandExecutor()
        self.memory = MemoryManager(config.DB_PATH)
        self.persistent_memory = PersistentMemory(config.PERSISTENT_MEMORY_DB)
        self.notes = NotesManager(config.NOTES_FILE)
        # Единый источник настроения (RAG передаём, чтобы не дублировать индекс)
        self.context = ContextManager(rag=self.rag)

        # Инициализация гибридного анализатора (лениво, если LAZY_MICRO_MODELS)
        self.analyzer = None
        self.emotional_analyzer = None
        if not getattr(config, "LAZY_MICRO_MODELS", True):
            self._init_analyzer()

        self.intent_router = None  # мёртвый путь, не создавать

        # Умный выбор анимаций (эмоции остаются)
        self.anim_selector = None
        if HAS_ANIM_SELECTOR and AnimationSelector:
            try:
                self.anim_selector = AnimationSelector(analyzer=self.analyzer)
                self.anim_selector.context = self.context
            except Exception as e:
                logger.warning(f"AnimationSelector: {e}")

        # Состояние
        self.conversation_history: List[Dict] = []
        self._lock = asyncio.Lock()
        self._initialized = False
        self._background_tasks = []
        self._reminder_callback = None
        self._system_alert_callback = None

        self.voice.set_callback(self._on_voice_input)
        self.server_ok = False
        try:
            from llm_server import ping
            self.server_ok = ping(1.2)
        except Exception:
            self.server_ok = False
        if getattr(config, "RAG_AUTO_INDEX", True) and getattr(config, "ADVANCED_RAG_ENABLED", True):
            if self.server_ok:
                asyncio.create_task(self._auto_index())
            else:
                logger.info("RAG автоиндексация отложена — LLM сервер не отвечает")
        else:
            logger.info("RAG автоиндексация пропущена (выключена флагами)")
        self._load_history()

        self._initialized = True
        logger.info(
            "LMAssistant инициализирован "
            f"(core=ok micro={getattr(config, 'USE_MICRO_MODELS', False)} "
            f"router=off "
            f"rag_auto={getattr(config, 'RAG_AUTO_INDEX', False)} "
            f"voice={getattr(config, 'ENABLE_VOICE_INPUT', False)} "
            f"greet={getattr(config, 'ENABLE_AUTO_GREETING', False)} "
            f"nsfw={getattr(config, 'NSFW_ENABLED', False)})"
        )

    def _init_analyzer(self):
        """Единый анализатор через emotion_service (без дублей в GUI)."""
        try:
            from emotion_service import get_shared_analyzer
            self.analyzer = get_shared_analyzer()
            self.emotional_analyzer = self.analyzer  # совместимость
            if self.analyzer:
                logger.info("✅ Анализатор эмоций/интентов: единый emotion_service")
            else:
                logger.warning("⚠️ emotion_service вернул None")
        except Exception as e:
            logger.error(f"❌ emotion_service: {e}")
            self.analyzer = None
            self.emotional_analyzer = None

    # ===== СОВМЕСТИМОСТЬ С GUI =====

    def init_voice(self) -> VoiceController:
        return self.voice

    def set_reminder_callback(self, callback):
        self._reminder_callback = callback
        if hasattr(self.executor, "reminder_manager"):
            self.executor.reminder_manager.set_callback(callback)

    def set_system_alert_callback(self, callback):
        self._system_alert_callback = callback
        if hasattr(self.executor, "set_alert_callback"):
            self.executor.set_alert_callback(callback)

    def check_reminders_now(self):
        if hasattr(self.executor, "reminder_manager"):
            self.executor.reminder_manager.check_reminders()

    def read_file_content(self, file_path: str, max_chars: int = 10000) -> str:
        import os
        ext = os.path.splitext(file_path)[1].lower()
        content = ""
        try:
            if ext in (
                ".txt", ".py", ".js", ".json", ".md", ".csv", ".log",
                ".xml", ".html", ".css", ".sql", ".yaml", ".yml", ".toml",
                ".ini", ".cfg", ".conf", ".sh", ".bat", ".ps1",
            ):
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
            elif ext == ".pdf":
                try:
                    import PyPDF2
                    with open(file_path, "rb") as f:
                        reader = PyPDF2.PdfReader(f)
                        for page in reader.pages:
                            content += page.extract_text() + "\n"
                except Exception as e:
                    content = f"Ошибка чтения PDF: {e}"
            elif ext == ".docx":
                try:
                    from docx import Document
                    doc = Document(file_path)
                    for para in doc.paragraphs:
                        content += para.text + "\n"
                except Exception as e:
                    content = f"Ошибка чтения DOCX: {e}"
            elif ext in (".xlsx", ".xls"):
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(file_path, data_only=True)
                    for sheet in wb.worksheets:
                        content += f"\n=== {sheet.title} ===\n"
                        for row in sheet.iter_rows(values_only=True):
                            content += " | ".join([str(c) if c is not None else "" for c in row]) + "\n"
                except Exception as e:
                    content = f"Ошибка чтения XLSX: {e}"
            elif ext == ".pptx":
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                except Exception as e:
                    content = f"Ошибка чтения PPTX: {e}"
            else:
                content = f"Тип файла {ext} не поддерживается."
        except Exception as e:
            content = f"Ошибка чтения файла: {e}"

        if len(content) > max_chars:
            content = content[:max_chars] + "\n\n... (обрезано)"
        return content

    def clear_history(self):
        self.conversation_history = []

    def _load_history(self):
        try:
            history = self.memory.get_recent_history(20, character=active_character())
            self.conversation_history = [
                msg for msg in history if msg["role"] in ("user", "assistant")
            ]
        except Exception as e:
            logger.error(f"Ошибка загрузки истории: {e}")

    def switch_character(self, new_name: str):
        """Смена персонажа = новый чат. Старую ленту не подмешиваем."""
        name = (new_name or "").strip() or "лисичка"
        old = active_character()
        config.ACTIVE_CHARACTER = name
        try:
            from character_manager import apply_to_config, apply_character_paths
            apply_character_paths(name)
            apply_to_config()
        except Exception as e:
            logger.warning(f"persona pack on switch: {e}")
        try:
            from memory_manager import MemoryManager
            from persistent_memory import PersistentMemory
            self.memory = MemoryManager(config.DB_PATH)
            self.persistent_memory = PersistentMemory(config.PERSISTENT_MEMORY_DB)
        except Exception as e:
            logger.warning(f"memory reopen: {e}")
        self.conversation_history = []
        if hasattr(self, "context") and self.context:
            try:
                self.context._user_mood = "neutral"
                self.context._user_mood_until = 0.0
            except Exception:
                pass
        logger.info(f"🎭 Персонаж {old} → {name}; история очищена")
        return []

    def _ensure_analyzer(self):
        if self.analyzer is None and self.emotional_analyzer is None:
            self._init_analyzer()
        return self.analyzer

    async def _auto_index(self):
        try:
            n = len(getattr(self.rag, "chunks", []) or [])
            min_chunks = int(getattr(config, "RAG_MIN_CHUNKS_FOR_SEMANTIC", 8) or 8)
            if n == 0:
                logger.info("RAG: индекс пуст — автоиндексация документов (без обязательных embeddings)")
            await self.rag.auto_index_from_config_async()
            n2 = len(getattr(self.rag, "chunks", []) or [])
            logger.info(f"RAG автоиндексация завершена (чанков: {n2})")
            if n2 < min_chunks and getattr(config, "RAG_LAZY_EMBEDDINGS", True):
                logger.info("RAG: мало чанков — semantic/embeddings не форсируем")
        except Exception as e:
            logger.error(f"Ошибка автоиндексации: {e}")

    # ===== ПРЕОБРАЗОВАНИЕ ИНТЕНТА В КОМАНДУ =====

    def _intent_to_command(self, intent: str, params: Dict) -> Optional[Dict]:
        """Преобразование интента в команду для CommandExecutor.
        Возвращает None, если обязательные параметры пустые —
        тогда запрос уходит в LLM, а не падает с «Не указано имя».
        """
        params = params or {}

        if intent == "search":
            q = (params.get("query") or "").strip()
            if not q:
                return None
            return {"type": "SEARCH", "params": q}

        if intent == "launch_app":
            app = (params.get("app") or "").strip()
            if not app:
                return None
            return {"type": "LAUNCH", "params": app}

        if intent == "open_browser":
            return {"type": "OPEN", "params": params.get("url") or "chrome"}

        if intent == "system_control":
            return self._system_control_to_command(params)

        if intent == "volume_control":
            return self._volume_control_to_command(params)

        if intent == "reminder":
            return {"type": "REMINDER", "params": params}

        if intent == "notes":
            text = (params.get("text") or "").strip()
            if not text:
                return None
            return {"type": "NOTE_ADD", "params": text}

        if intent == "screenshot":
            return {"type": "SCREENSHOT", "params": None}

        if intent == "file_operation":
            path = (params.get("path") or "").strip()
            if not path:
                return None
            return {"type": "OPEN", "params": path}

        return None

    def _system_control_to_command(self, params: Dict) -> Optional[Dict]:
        """Преобразование системного контроля в команду."""
        action = params.get('action', '')
        if action == 'shutdown':
            return {'type': 'SHUTDOWN', 'params': {'confirm': False}}
        elif action == 'restart':
            return {'type': 'RESTART', 'params': {'confirm': False}}
        elif action == 'lock':
            return {'type': 'LOCK', 'params': None}
        elif action == 'minimize_all':
            return {'type': 'MINIMIZE', 'params': 'all'}
        elif action == 'maximize_all':
            return {'type': 'MAXIMIZE', 'params': 'all'}
        return None

    def _volume_control_to_command(self, params: Dict) -> Optional[Dict]:
        """Преобразование управления громкостью в команду."""
        action = params.get('action', '')
        if action == 'up':
            return {'type': 'VOLUME_UP', 'params': None}
        elif action == 'down':
            return {'type': 'VOLUME_DOWN', 'params': None}
        elif action == 'mute':
            return {'type': 'MUTE', 'params': None}
        elif 'level' in params:
            return {'type': 'VOLUME', 'params': params.get('level')}
        return None

    # ===== ОСНОВНАЯ ЛОГИКА =====

    async def generate_stream(
        self,
        user_message: str,
        stream_callback: Optional[Callable[[str], None]] = None,
        image_path: Optional[str] = None,
        file_content: Optional[str] = None,
    ) -> AsyncGenerator[str, None]:
        if not user_message and not file_content:
            try:
                from character_manager import character_short_name
                yield "Скажи что-нибудь — я %s." % character_short_name()
            except Exception:
                yield "Скажи что-нибудь."
            return

        if file_content:
            user_message = f"{user_message}\n\nСодержимое файла:\n{file_content[:5000]}"

        async with self._lock:
            self.conversation_history.append({"role": "user", "content": user_message})
            self.memory.add_message("user", user_message, character=active_character())
        self._maybe_store_user_facts(user_message)

        # Активность + настроение (единый источник)
        self.context.touch_activity()
        self.context.detect_mood(user_message)

        # Один маршрутизатор: явная команда → исполнить → 1–2 фразы → стоп.
        local = None
        try:
            from commands import parse_user, execute, flavor, needs_scanner, handle_user
            cmd = parse_user(user_message)
            if cmd and needs_scanner(cmd) and getattr(self, "executor", None):
                scanned = await self.executor.execute_async(
                    {"type": "LAUNCH", "params": cmd.target}
                )
                if scanned:
                    line = flavor(cmd.kind, scanned)
                    local = f"[ANIM:{cmd.anim}] {line}\n{scanned}"
            elif cmd:
                local = handle_user(user_message)
        except Exception as e:
            logger.error(f"commands.handle_user: {e}")
            local = None
        if local:
            yield local
            hist = self._clean_for_history(local)
            if not hist:
                log_system("cmd", local)
                hist = clean_reply(local, 120) or "готово"
            async with self._lock:
                self.conversation_history.append({"role": "assistant", "content": hist})
                self.memory.add_message("assistant", hist, character=active_character())
            return

        # Нет сервера — GUI и команды уже живы, в модель не ходим и не виснем.
        try:
            from llm_server import ping, offline_reply
        except Exception:
            ping = lambda *a, **k: False
            def offline_reply():
                return "[ANIM:thinking] Сервер ИИ выключен."
        if not ping():
            yield offline_reply()
            return
        self.server_ok = True

        # Эмоции / морда — не команды. Анализатор только для AnimationSelector.
        analyzer = self._ensure_analyzer()
        if self.anim_selector and analyzer and getattr(self.anim_selector, "analyzer", None) is None:
            self.anim_selector.analyzer = analyzer

        # ===== УМНЫЙ ВЫБОР АНИМАЦИИ =====
        anim = "neutral"
        if getattr(self, "anim_selector", None):
            try:
                anim = self.anim_selector.select(
                    user_text=user_message,
                )
            except Exception as e:
                logger.error(f"AnimationSelector: {e}")
        elif self.analyzer:
            try:
                anim = self.analyzer.get_animation(user_message)
            except Exception as e:
                logger.error(f"Ошибка анализа эмоций: {e}")
        elif hasattr(self, "emotional_analyzer") and self.emotional_analyzer:
            try:
                emotion, _ = self.emotional_analyzer.analyze_full_context(user_message)
                anim = (
                    self.emotional_analyzer.get_animation(emotion)
                    if hasattr(self.emotional_analyzer, "get_animation")
                    else "neutral"
                )
            except Exception as e:
                logger.error(f"Ошибка rule-based анализа эмоций: {e}")

        # ===== ПАРАЛЛЕЛЬНЫЙ СБОР КОНТЕКСТА =====
        # RAG + память с жёстким лимитом ожидания (не блокируем UI на 15+ сек)
        rag_tasks = []
        if getattr(config, "ADVANCED_RAG_ENABLED", False) and getattr(self, "rag", None):
            mode = (getattr(config, "RAG_DEFAULT_MODE", "hybrid") or "hybrid")
            if getattr(config, "RAG_HYBRID_MODE", False):
                mode = "hybrid"
            rag_tasks.append(
                self.rag.get_context_async(
                    user_message,
                    limit=int(getattr(config, "RAG_CONTEXT_LIMIT", 4) or 4),
                    max_chars=int(getattr(config, "RAG_MAX_CONTEXT_CHARS", 2200) or 2200),
                    mode=mode,
                    min_similarity=float(getattr(config, "RAG_MIN_SIMILARITY", 0.22) or 0.22),
                )
            )
        memory_task = asyncio.create_task(self._get_memory_context_async(user_message))
        rag_tasks.append(memory_task)

        context_parts = []
        try:
            rag_results = await asyncio.wait_for(
                asyncio.gather(*rag_tasks, return_exceptions=True),
                timeout=float(getattr(config, "RAG_CONTEXT_TIMEOUT", 3.0)),
            )
            for result in rag_results:
                if isinstance(result, str) and result:
                    context_parts.append(result)
                elif isinstance(result, Exception):
                    logger.error(f"Ошибка RAG/памяти: {result}")
        except asyncio.TimeoutError:
            logger.warning("RAG/память: таймаут 3с — отвечаем без контекста")
            # отменяем зависшие
            for t in rag_tasks:
                if isinstance(t, asyncio.Task) and not t.done():
                    t.cancel()

        context = "\n".join(context_parts) if context_parts else ""
        mood_addon = self.context.get_mood_prompt_addon()
        if self._is_find_images_request(user_message) and not self._is_image_prompt_request(user_message):
            mood_addon = (mood_addon or "") + (
                "\nХозяин просит НАЙТИ готовые картинки в интернете. "
                "Коротко подтверди и напиши [SEARCH картинки …]. "
                "НЕ пиши промпты для генерации и не выдумывай ссылки."
            )
        elif self._is_image_prompt_request(user_message):
            mood_addon = (mood_addon or "") + (
                "\nХозяин просит именно ПРОМПТЫ для генерации картинки, не поиск в браузере."
            )

        screen_jpeg = None
        look_screen = False
        try:
            from screen_watch import (
                is_look_command,
                capture_jpeg,
                user_content_with_image,
                VISION_ADDON,
                extract_scene,
            )
            look_screen = bool(image_path) or (
                getattr(config, "SCREEN_VISION_ENABLED", True)
                and is_look_command(user_message)
            )
            if look_screen and not image_path:
                hide = getattr(self, "hide_for_screenshot", None)
                show = getattr(self, "show_after_screenshot", None)
                try:
                    if callable(hide):
                        hide()
                    screen_jpeg = capture_jpeg(text=user_message)
                finally:
                    if callable(show):
                        show()
                if screen_jpeg:
                    image_path = screen_jpeg
                    focus = ""
                    try:
                        meta = os.path.join(
                            getattr(config, "DATA_DIR", "."), "cache", "screen_last.txt"
                        )
                        if os.path.isfile(meta):
                            with open(meta, encoding="utf-8") as _mf:
                                focus = _mf.read().strip()
                    except Exception:
                        focus = ""
                    extra = f"\nСейчас во вложении: {focus}." if focus else ""
                    ocr_block = ""
                    try:
                        from screen_watch import last_ocr_text
                        ocr = last_ocr_text()
                        if ocr:
                            ocr_block = (
                                "\n\n[OCR текста с выбранного монитора — подсказка для мелкого UI]\n"
                                + ocr
                            )
                    except Exception:
                        ocr_block = ""
                    mood_addon = (mood_addon or "") + VISION_ADDON + extra + ocr_block
                    anim = "searching"
            elif image_path and look_screen:
                mood_addon = (mood_addon or "") + VISION_ADDON
        except Exception as e:
            logger.error(f"screen_watch: {e}")

        # ===== МОДУЛЬНЫЙ SYSTEM PROMPT =====
        if HAS_PROMPT_BUILDER and build_system_prompt:
            system_prompt = build_system_prompt(
                mood_addon=mood_addon or "",
                rag_context=context,
                memory_context="",
                include_few_shot=len(self.conversation_history) < 3,
                include_commands=len(self.conversation_history) < 3,
                include_anim_hints=len(self.conversation_history) < 3,
                max_context_chars=int(getattr(config, "RAG_MAX_CONTEXT_CHARS", 1200) or 1200),
            )
        else:
            system_prompt = self.system_prompt
            try:
                from time_context import get_time_prompt_addon
                _tb = get_time_prompt_addon()
                if _tb:
                    system_prompt += _tb
            except Exception:
                pass
            if context:
                system_prompt += f"\n\n{context}"
            if mood_addon:
                system_prompt += mood_addon


        try:
            from character_manager import build_character_prompt_block
            card = build_character_prompt_block()
            if card:
                system_prompt = card + "\n\n" + system_prompt
        except Exception as e:
            logger.warning(f"character card prompt: {e}")

        try:
            from files_box import prompt_block
            box = prompt_block()
            if box:
                system_prompt += "\n\n" + box
        except Exception:
            pass

        messages = [{"role": "system", "content": system_prompt}]

        history_limit = int(getattr(config, "CHAT_HISTORY_TURNS", 6) or 6)
        hist_chars = int(getattr(config, "CHAT_HISTORY_CHARS", 360) or 360)
        for msg in self.conversation_history[-history_limit:-1]:
            role = msg.get("role") or "user"
            raw = msg.get("content") or ""
            if not isinstance(raw, str):
                continue
            content = self._clean_for_history(raw, limit=hist_chars)
            if not content:
                continue
            messages.append({"role": role, "content": content})

        user_cap = int(getattr(config, "CHAT_USER_CHARS", 800) or 800)
        user_payload = user_message
        if isinstance(user_payload, str) and len(user_payload) > user_cap:
            user_payload = user_payload[:user_cap].rsplit(" ", 1)[0] + "…"
        if image_path:
            try:
                from screen_watch import user_content_with_image
                user_payload = user_content_with_image(
                    user_payload if isinstance(user_payload, str) else user_message,
                    image_path,
                )
            except Exception as e:
                logger.error(f"vision payload: {e}")
                user_payload = user_message
        messages.append({"role": "user", "content": user_payload})

        async with api_semaphore:
            async with aiohttp.ClientSession() as session:
                use_model = self.model_name
                if image_path or look_screen:
                    use_model = (
                        getattr(config, "SCREEN_VISION_MODEL", None)
                        or getattr(config, "FAST_MODEL", None)
                        or self.model_name
                    )
                payload = {
                    "model": use_model,
                    "messages": messages,
                    "temperature": self.temperature,
                    "max_tokens": self.max_tokens,
                    "stream": True,
                }
                try:
                    async with session.post(
                        f"{self.api_url}/chat/completions",
                        headers={"Authorization": f"Bearer {self.api_key}"},
                        json=payload,
                        timeout=aiohttp.ClientTimeout(
                            total=float(getattr(config, "LLM_TIMEOUT", 300)),
                            sock_connect=float(getattr(config, "LLM_CONNECT_TIMEOUT", 3)),
                            sock_read=float(getattr(config, "LLM_SOCK_READ_TIMEOUT", 300)),
                        ),
                    ) as response:
                        if response.status != 200:
                            error_text = await response.text()
                            yield f"Ошибка API: {response.status} - {error_text}"
                            return

                        full_response = ""
                        first_chunk = True
                        async for line in response.content:
                            if not line:
                                continue
                            try:
                                line = line.decode("utf-8", errors="replace")
                            except Exception:
                                continue

                            if line.startswith("data: "):
                                data_str = line[6:].strip()
                                if data_str == "[DONE]":
                                    break
                                try:
                                    data = json.loads(data_str)
                                    delta = data.get("choices", [{}])[0].get("delta", {})
                                    chunk = delta.get("content", "")
                                    if chunk:
                                        full_response += chunk
                                        if first_chunk:
                                            # На скрине ждём [ANIM] от самой модели
                                            if look_screen:
                                                first_chunk = False
                                            elif not re.search(r'\[ANIM:\w+\]', full_response, re.IGNORECASE):
                                                chunk = f"[ANIM:{anim}] {chunk}"
                                                full_response = chunk
                                                first_chunk = False
                                        yield chunk
                                        if stream_callback:
                                            stream_callback(chunk)
                                except json.JSONDecodeError:
                                    continue

                        if full_response:
                            hist = self._clean_for_history(full_response)
                            async with self._lock:
                                self.conversation_history.append(
                                    {"role": "assistant", "content": hist or full_response[:400]}
                                )
                                self.memory.add_message(
                                    "assistant",
                                    hist or full_response[:400],
                                    character=active_character(),
                                )

                            # ===== ОБРАБОТКА КОМАНД ПАМЯТИ =====
                            await self._process_memory_commands(full_response)

                            # Команды — ТОЛЬКО через executor
                            user_q = self._extract_user_search_query(
                                user_message,
                                previous=self._previous_user_text(),
                            )
                            find_images = self._is_find_images_request(user_message)
                            want_prompt = self._is_image_prompt_request(user_message)
                            search_done = False

                            if CommandParser.has_commands(full_response):
                                commands = CommandParser.parse(full_response)
                                for cmd in commands:
                                    if cmd["type"] == "ANIM":
                                        continue
                                    raw_p = str(cmd.get("params") or "").strip().lower()
                                    if look_screen and cmd["type"] in (
                                        "SEARCH", "LAUNCH", "OPEN", "RUN"
                                    ):
                                        logger.info(
                                            f"пропуск {cmd['type']}: запрос был просмотр экрана"
                                        )
                                        continue
                                    if raw_p in (
                                        "путь", "path", "название", "команда",
                                        "запрос", "query", "файл", "текст", "none",
                                        "что-нибудь", "что нибудь", "ничего",
                                        "всё", "все", "это", "то",
                                    ):
                                        logger.info(f"пропуск {cmd['type']}: шаблонный параметр {raw_p!r}")
                                        continue
                                    # «найди картинки» ≠ открыть браузер и ≠ промпт
                                    if find_images and not want_prompt and cmd["type"] in (
                                        "OPEN", "LAUNCH", "RUN"
                                    ):
                                        logger.info(f"пропуск {cmd['type']}: запрос был поиск картинок")
                                        continue
                                    if self._looks_like_launch(user_message) and cmd["type"] == "SEARCH":
                                        logger.info("пропуск SEARCH: запрос был запуск программы")
                                        continue
                                    if cmd["type"] == "SEARCH" and user_q:
                                        cmd = dict(cmd)
                                        q = user_q
                                        if find_images and not want_prompt and not re.search(
                                            r"картин|фото|image", q, re.I
                                        ):
                                            q = f"картинки {q}"
                                        cmd["params"] = q
                                    result = await self.executor.execute_async(cmd)
                                    if result:
                                        log_system(cmd.get("type") or "cmd", result)
                                    if (
                                        cmd["type"] == "SEARCH"
                                        and result
                                        and "пропущен" not in (result or "")
                                    ):
                                        search_done = True

                            launch_target = self._extract_launch_target(user_message)
                            launched = any(
                                c.get("type") in ("LAUNCH", "OPEN", "NOTEPAD")
                                for c in (
                                    CommandParser.parse(full_response)
                                    if CommandParser.has_commands(full_response)
                                    else []
                                )
                            )
                            if (
                                not look_screen
                                and self._looks_like_launch(user_message)
                                and launch_target
                                and not launched
                            ):
                                result = await self.executor.execute_async(
                                    {"type": "LAUNCH", "params": launch_target}
                                )
                                if result:
                                    log_system("LAUNCH", result)
                                launched = True

                            # Второй поиск отключён: явный «найди» уже ушёл в handle_user.
                            if user_q and not search_done and not look_screen:
                                logger.info(
                                    f"🔍 search tag пропущен (нет явной команды в этом ходе): {user_q!r}"
                                )

                except asyncio.TimeoutError:
                    yield "⏰ Таймаут соединения с API"
                except aiohttp.ClientError as e:
                    yield f"🌐 Ошибка соединения: {e}"
                except Exception as e:
                    logger.error(f"Ошибка генерации: {e}")
                    yield f"❌ Ошибка: {e}"

    async def send_message_async(
        self,
        user_message: str,
        stream_callback: Optional[Callable[[str], None]] = None,
        image_path: Optional[str] = None,
        file_content: Optional[str] = None,
    ) -> str:
        full_response = ""
        async for chunk in self.generate_stream(
            user_message, stream_callback, image_path, file_content
        ):
            full_response += chunk
        return full_response

    # ===== ПАМЯТЬ =====

    async def _get_memory_context_async(self, query: str) -> str:
        try:
            context = self.persistent_memory.get_context_for_prompt(
                query=query,
                scope=prompt_scopes(),
                limit=config.MAX_MEMORIES_IN_CONTEXT,
            )
            return context or ""
        except Exception as e:
            logger.error(f"Ошибка памяти: {e}")
            return ""

    async def _process_memory_commands(self, reply: str):
        """
        Обрабатывает команды памяти из ответа ассистента.
        Поддерживает:
        - [REMEMBER категория: значение]
        - [REMEMBER_APP имя путь]
        - [REMEMBER_ALIAS алиас цель]
        - [REMEMBER_PROJECT имя: описание]
        
        Особое внимание: [REMEMBER app: имя путь] и [REMEMBER_APP имя путь]
        автоматически добавляют программу в AppScanner и алиасы.
        """
        if not self.persistent_memory:
            return
        
        # ===== ОБРАБОТКА REMEMBER (ОБЫЧНЫЙ) =====
        for match in re.finditer(
            r"\[REMEMBER\s+(.+?):\s*(.+?)\]", reply, re.IGNORECASE
        ):
            category = match.group(1).strip().lower()
            value = match.group(2).strip()
            if is_system_or_ocr(value):
                log_system("remember_skip", value)
                continue
            key = f"user_{category}_{value[:30]}"
            
            # Сохраняем в память
            self.persistent_memory.add_memory(
                scope_for_category(category), category, key, value, confidence=0.8
            )
            logger.info(f"Запомнено: {category} → {value[:50]}")
            
            # ===== ЕСЛИ ЭТО ПРОГРАММА — ДОБАВЛЯЕМ В АЛИАСЫ =====
            if category in ("app", "программа", "program", "софт", "приложение"):
                # Пробуем разобрать "имя путь"
                # Ищем путь: либо в кавычках, либо всё что после пробела
                
                # Сначала ищем путь в кавычках
                quoted = re.search(r'"([^"]+)"', value)
                if quoted:
                    path = quoted.group(1)
                    # Имя — всё до кавычек
                    name = value[:quoted.start()].strip()
                else:
                    # Ищем по расширению .exe или .lnk
                    exe_match = re.search(r'(.+?)\s+(.+?\.(?:exe|lnk))', value, re.I)
                    if exe_match:
                        name = exe_match.group(1).strip()
                        path = exe_match.group(2).strip()
                    else:
                        # Пробуем разделить по пробелу
                        parts = value.split(maxsplit=1)
                        if len(parts) >= 2:
                            name = parts[0].strip()
                            path = parts[1].strip()
                        else:
                            # Если нет пробела — возможно это уже имя алиаса
                            name = value
                            # Пробуем найти существующий алиас
                            if self.executor.app_scanner:
                                existing = self.executor.app_scanner.resolve_alias(value)
                                if existing:
                                    logger.info(f"Алиас уже существует: {value} → {existing['target']}")
                                    continue
                            else:
                                logger.info(f"Запомнено как текст: {category} → {value[:50]}")
                                continue
                
                # Если путь существует или похож на программу
                if path and (os.path.exists(path) or path.lower().endswith(('.exe', '.lnk'))):
                    if self.executor.app_scanner:
                        # Добавляем в apps
                        self.executor.app_scanner.add_manual_app(name, path, "alias")
                        # Добавляем алиас
                        self.executor.app_scanner.add_alias(name, path, "app")
                        logger.info(f"✅ Программа добавлена в алиасы: {name} → {path}")
                        
                        # Обновляем persistent_memory с правильным форматом
                        self.persistent_memory.add_memory(
                            "pc", "app_alias", name, path, confidence=0.95
                        )
                else:
                    # Если это не программа — пробуем найти существующий алиас
                    if self.executor.app_scanner:
                        alias_data = self.executor.app_scanner.resolve_alias(value)
                        if alias_data:
                            logger.info(f"Алиас уже существует: {value} → {alias_data['target']}")
                        else:
                            logger.info(f"Запомнено как текст: {category} → {value[:50]}")

        # ===== ОБРАБОТКА REMEMBER_APP (СТАРЫЙ ФОРМАТ) =====
        for match in re.finditer(
            r"\[REMEMBER_APP\s+(.+?)\s+(.+?)\]", reply, re.IGNORECASE
        ):
            name = match.group(1).strip()
            path = match.group(2).strip()
            if self.executor.app_scanner:
                self.executor.app_scanner.add_manual_app(name, path, "alias")
                self.executor.app_scanner.add_alias(name, path, "app")
                self.persistent_memory.add_memory(
                    "pc", "app_alias", name, path, confidence=0.95
                )
                logger.info(f"✅ Запомнена программа (REMEMBER_APP): {name} → {path}")

        # ===== ОБРАБОТКА REMEMBER_ALIAS =====
        for match in re.finditer(
            r"\[REMEMBER_ALIAS\s+(.+?)\s+(.+?)(?:\s+as\s+(\w+))?\]", reply, re.IGNORECASE
        ):
            alias = match.group(1).strip()
            target = match.group(2).strip()
            type_ = (match.group(3) if match.group(3) else "").strip().lower()
            if self.executor.app_scanner:
                # Автоопределение типа
                if not type_ or type_ not in ('app', 'file', 'folder', 'url', 'command'):
                    if target.startswith(('http://', 'https://', 'www.')):
                        type_ = 'url'
                    elif os.path.exists(target):
                        if os.path.isdir(target):
                            type_ = 'folder'
                        elif target.lower().endswith('.exe') or target.lower().endswith('.lnk'):
                            type_ = 'app'
                        else:
                            type_ = 'file'
                    else:
                        type_ = 'command'
                
                if type_ == 'app' and (target.lower().endswith('.exe') or target.lower().endswith('.lnk')):
                    self.executor.app_scanner.add_manual_app(alias, target, "alias")
                self.executor.app_scanner.add_alias(alias, target, type_)
                self.persistent_memory.add_memory(
                    "pc", "app_alias", alias, target, confidence=0.95
                )
                logger.info(f"✅ Запомнен алиас (REMEMBER_ALIAS): {alias} → {target} ({type_})")

        # ===== ОБРАБОТКА ПРОЕКТОВ =====
        for match in re.finditer(
            r"\[REMEMBER_PROJECT\s+(.+?):\s*(.+?)\]", reply, re.IGNORECASE
        ):
            name = match.group(1).strip()
            description = match.group(2).strip()
            self.persistent_memory.add_memory(
                "project", "project", name, description, confidence=0.9
            )
            logger.info(f"Запомнен проект: {name} → {description[:50]}")

    # ===== ПОИСКОВЫЙ ЗАПРОС ИЗ СООБЩЕНИЯ ХОЗЯИНА =====

    @staticmethod
    def _clean_for_history(text: str, limit: int = 180) -> str:
        """В историю и в модель — короткая реплика. Система/OCR не пишем сюда."""
        return clean_reply(text, limit=limit)

    def _maybe_store_user_facts(self, user_message: str):
        """Только явные факты, не весь чат."""
        if not user_message or not getattr(self, "persistent_memory", None):
            return
        t = user_message.strip()
        if is_system_or_ocr(t):
            log_system("user_skip", t)
            return
        pairs = []
        m = re.search(r"(?:запомни(?:\s+что)?|remember that)\s+(.+)$", t, flags=re.I)
        if m:
            val = m.group(1).strip()[:400]
            if val:
                pairs.append(("fact", val[:60], val, 0.85))
        m = re.search(r"меня зовут\s+([A-Za-zА-Яа-яЁё\- ]{2,40})", t, flags=re.I)
        if m:
            name = m.group(1).strip()
            pairs.append(("profile", "name", name, 0.95))
        m = re.search(r"я\s+(?:живу|из)\s+(.+)$", t, flags=re.I)
        if m:
            pairs.append(("profile", "city", m.group(1).strip()[:80], 0.7))
        for cat, key, val, conf in pairs:
            try:
                self.persistent_memory.add_memory(
                    scope_for_category(cat), cat, key, val, confidence=conf, importance=conf
                )
                logger.info(f"Persistent fact: {cat}/{key}={val[:50]}")
            except Exception as e:
                logger.warning(f"fact store: {e}")

    def _previous_user_text(self) -> str:
        try:
            for msg in reversed(self.conversation_history[:-1]):
                if msg.get("role") == "user":
                    return str(msg.get("content") or "")
        except Exception:
            pass
        return ""

    @staticmethod
    def _looks_like_launch(text: str) -> bool:
        t = (text or "").lower()
        if not re.search(r"\b(запусти|открой|открыть|запуск)\b", t):
            return False
        if any(w in t for w in ("картин", "фото", "сайт", "вкладк", "гугл", "яндекс", "http")):
            return False
        return True

    @staticmethod
    def _extract_launch_target(text: str) -> Optional[str]:
        t = (text or "").strip()
        m = re.search(
            r"(?:запусти|открой|открыть|запуск)\s+(.+)$",
            t,
            flags=re.I,
        )
        if not m:
            return None
        app = m.group(1).strip(" .!?,")
        app = re.sub(r"^(пожалуйста|плиз|мне)\s+", "", app, flags=re.I)
        aliases = {
            "блокнот": "блокнот",
            "блакнот": "блокнот",
            "notepad": "блокнот",
            "калькулятор": "калькулятор",
            "калкулятор": "калькулятор",
            "проводник": "проводник",
        }
        key = app.lower()
        return aliases.get(key, app) if len(app) >= 2 else None

    @staticmethod
    def _is_image_prompt_request(text: str) -> bool:
        t = (text or "").lower()
        keys = (
            "промпт", "prompt", "для генерац", "сгенерируй описание",
            "описание для картин", "для midjourney", "для stable",
            "для flux", "negative prompt",
        )
        return any(k in t for k in keys)

    @staticmethod
    def _is_find_images_request(text: str) -> bool:
        t = (text or "").lower()
        if LMAssistant._is_image_prompt_request(t):
            return False
        img = ("картин", "кортин", "фото", "изображен", "арт", "pic", "image")
        find = ("найд", "нади", "поищ", "поиск", "покажи", "скинь", "гугл", "найди")
        return any(w in t for w in img) and any(w in t for w in find)

    @staticmethod
    def _extract_user_search_query(
        user_message: str, previous: str = ""
    ) -> Optional[str]:
        if not user_message or not isinstance(user_message, str):
            return None
        text = user_message.strip()
        triggers = (
            r"(?:найд[иуёе]?|нади|поищи|поиск|гугл|google|find|search|look\s*up)"
        )
        t = re.sub(
            r"^(?:лисичка|лис[ая]|пожалуйста|плиз|hey|ok)[,\s]+",
            "",
            text,
            flags=re.I,
        )
        m = re.search(triggers + r"\s+(.+)$", t, flags=re.IGNORECASE | re.DOTALL)
        if not m:
            m2 = re.search(
                r"(?:картинк\w*|кортинк\w*|фото|изображен\w*|арт|pics?|images?)\s+(.+)$",
                t,
                flags=re.IGNORECASE | re.DOTALL,
            )
            if not m2:
                return None
            q = m2.group(1).strip()
        else:
            q = m.group(1).strip()
        q = re.sub(r"[.!?…]+$", "", q).strip()
        q = q.strip("«»\"'“”").strip()
        q = re.sub(r"\s+", " ", q)
        if len(q) < 2:
            return None
        if q.lower() in ("уже выполнен", "(уже выполнен)", "none", "null"):
            return None
        return q[:200]

    # ===== ГОЛОС =====

    async def _on_voice_input(self, text: str):
        if text:
            logger.info(f"Голос: {text[:50]}")
            response = await self.send_message_async(text)
            from speech_text import for_speech
            said = for_speech(response or "")
            if said:
                await self.voice.speak_async(said)

    async def process_voice_command_async(self) -> Optional[str]:
        text = await self.voice.listen_async()
        if text:
            response = await self.send_message_async(text)
            from speech_text import for_speech
            said = for_speech(response or "")
            if said:
                await self.voice.speak_async(said)
            return response
        return None

    # ===== УПРАВЛЕНИЕ =====

    async def clear_history_async(self):
        async with self._lock:
            self.conversation_history = []
            logger.info("История очищена")

    def shutdown(self):
        """Синхронная остановка — можно звать из closeEvent / aboutToQuit."""
        if getattr(self, "_shutting_down", False):
            return
        self._shutting_down = True
        logger.info("Shutdown...")
        try:
            lc = getattr(self, "_lifecycle", None) or getattr(self, "lifecycle", None)
            if lc is not None and hasattr(lc, "stop"):
                lc.stop()
        except Exception as e:
            logger.warning(f"Lifecycle stop: {e}")
        voice = getattr(self, "voice", None)
        if voice is not None:
            for meth in ("stop_listening", "stop_speaking", "close"):
                try:
                    fn = getattr(voice, meth, None)
                    if callable(fn):
                        fn()
                except Exception:
                    pass
        for name, obj in (
            ("memory", getattr(self, "memory", None)),
            ("persistent_memory", getattr(self, "persistent_memory", None)),
            ("rag", getattr(self, "rag", None)),
            ("context", getattr(self, "context", None)),
            ("executor", getattr(self, "executor", None)),
        ):
            try:
                if obj is not None and hasattr(obj, "close"):
                    obj.close()
            except Exception as e:
                logger.warning(f"Ошибка закрытия {name}: {e}")
        self.analyzer = None
        self.emotional_analyzer = None
        self.anim_selector = None
        try:
            import gc
            gc.collect()
        except Exception:
            pass
        logger.info("Shutdown done")

    async def shutdown_async(self):
        """Если цикл ещё жив — доп. остановка голоса. Иначе только shutdown()."""
        self.shutdown()
        try:
            asyncio.get_running_loop()
        except RuntimeError:
            return
        voice = getattr(self, "voice", None)
        if voice is None:
            return
        try:
            if hasattr(voice, "stop_listening_async"):
                await voice.stop_listening_async()
        except Exception as e:
            logger.warning(f"Ошибка остановки прослушивания: {e}")
        try:
            if hasattr(voice, "stop_speaking_async"):
                await voice.stop_speaking_async()
        except Exception as e:
            logger.warning(f"Ошибка остановки речи: {e}")

    def close(self):
        """Sync-обёртка (на крайний случай)."""
        try:
            lc = getattr(self, "_lifecycle", None)
            if lc is not None and hasattr(lc, "stop"):
                lc.stop()
        except Exception:
            pass
        try:
            if hasattr(self, "voice") and self.voice:
                self.voice.close()
        except Exception:
            pass
        for attr in ("memory", "persistent_memory", "rag", "context", "executor"):
            try:
                obj = getattr(self, attr, None)
                if obj is not None and hasattr(obj, "close"):
                    obj.close()
            except Exception:
                pass